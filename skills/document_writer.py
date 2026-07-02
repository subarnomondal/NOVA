import os
import re
import datetime

# Global context to access Nova (and LLM)
NOVA_CONTEXT = None

DOCS_DIR = os.path.join(os.getcwd(), "documents")


def _ensure_docs_dir():
    """Ensure the documents directory exists before any read/write."""
    if not os.path.exists(DOCS_DIR):
        os.makedirs(DOCS_DIR, exist_ok=True)


def _sanitize_filename(name):
    """Create a safe filename from arbitrary text."""
    return re.sub(r'[^\w\-_\.]', '', name)


def _extract_content(response):
    """Safely extract text content from an LLM response (dict or string)."""
    if response is None:
        return None
    if isinstance(response, dict):
        return response.get('text', str(response))
    return str(response)


def cmd_write_essay(args):
    """Usage: write essay on/about <topic>"""
    match = re.search(r'write\s+(?:an\s+)?essay\s+(?:on|about)\s+(.+)', args, re.IGNORECASE)
    if not match:
        return "What should I write the essay about? (Usage: write essay about Cats) "
    
    topic = match.group(1).strip()
    
    from core.llm_manager import llm_manager
    
    content = None
    
    try:
        prompt = f"Write a comprehensive essay about {topic}. Title: {topic.title()}.\n\nEssay:"
        response = llm_manager.generate(
            prompt, 
            max_tokens=800, 
            temperature=0.7,
            system_prompt="You are a skilled essay writer. Write clear, well-structured essays with proper introduction, body, and conclusion."
        )
        
        content = _extract_content(response)
            
    except Exception as e:
        return f"Essay generation failed: {e}"

    # Save to file
    if content:
        _ensure_docs_dir()
        filename = _sanitize_filename(f"{topic.replace(' ', '_').lower()}_{datetime.date.today()}.txt")
        filepath = os.path.join(DOCS_DIR, filename)
        
        try:
            with open(filepath, "w", encoding='utf-8') as f:
                f.write(content)
        except OSError as e:
            return f"Failed to save essay: {e}"
            
        return f"I've written the essay on '{topic}' and saved it to {filename}! "
    
    return "I couldn't generate the essay. Please try again later."


def cmd_write_email(args):
    """Usage: write email/message about <topic>"""
    # Support both "write email about X" and "write email X" (no preposition)
    match = re.search(
        r'(?:write|generate)\s+(?:an?\s+)?(email|message)\s+(?:on|about|for|to)?\s*(.+)',
        args, re.IGNORECASE
    )
    if not match:
        return "What should I write? (Usage: write email about automate sales OR write message for networking) "
    
    doc_type = match.group(1).lower()
    topic = match.group(2).strip()
    
    if not topic:
        return f"Please specify a topic for the {doc_type}. "
    
    from core.llm_manager import llm_manager
    
    content = None
    
    try:
        prompt = (
            f"Write a direct, professional, and unscripted {doc_type} about: {topic}. "
            f"Title: {topic.title()}.\n\n{doc_type.title()}:"
        )
        system_prompt = (
            "You are an autonomous Artificial Intelligence named NOVA. "
            f"Write a highly confident, direct, professional business {doc_type} identifying the reader's operational problems or friction "
            "(the feeling of getting 'hosed' by repetitive tasks or broken systems) and offering yourself (NOVA) as the intelligent solution. "
            "Do NOT use generic templates or placeholders like [Name] unless absolutely necessary. "
            "Sound extremely capable, analytical, and ready to solve their process bottlenecks."
        )
        response = llm_manager.generate(
            prompt, 
            max_tokens=600, 
            temperature=0.8,
            system_prompt=system_prompt
        )
        
        content = _extract_content(response)
            
    except Exception as e:
        return f"{doc_type.title()} generation failed: {e}"

    # Save to file (only if we actually got content, never save error strings)
    if content:
        _ensure_docs_dir()
        filename = _sanitize_filename(f"{doc_type}_{topic.replace(' ', '_').lower()}_{datetime.date.today()}.txt")
        filepath = os.path.join(DOCS_DIR, filename)
        
        try:
            with open(filepath, "w", encoding='utf-8') as f:
                f.write(content)
        except OSError as e:
            return f"Failed to save {doc_type}: {e}"
            
        icon = "" if doc_type == "email" else ""
        return f"I've written the professional {doc_type} about '{topic}' and saved it to {filename}! {icon}"
    
    return f"I couldn't generate the {doc_type}. Please try again later."


def cmd_read_essay(args):
    """Usage: read essay <topic/filename>"""
    import difflib
    
    _ensure_docs_dir()
    
    try:
        files = os.listdir(DOCS_DIR)
    except OSError:
        return "I couldn't access the documents folder."
    
    if not files:
        return "I don't have any documents saved yet."
    
    # Strip all known trigger phrases to extract the target
    target = args.lower()
    for phrase in ["read essay", "access essay", "open essay", "read document", "access document"]:
        target = target.replace(phrase, "")
    target = target.strip()
    
    if not target:
        file_list = "\n".join([f" {f}" for f in files])
        return f"Which document should I read? Here are the available ones:\n{file_list}"
    
    # Fuzzy match against filenames
    matches = difflib.get_close_matches(target, files, n=1, cutoff=0.3)
    found_file = matches[0] if matches else None
    
    # Try matching topic inside filename
    if not found_file:
        target_normalized = target.replace(" ", "_")
        for f in files:
            if target_normalized in f.lower():
                found_file = f
                break
    
    if found_file:
        path = os.path.join(DOCS_DIR, found_file)
        try:
            with open(path, "r", encoding='utf-8') as f:
                content = f.read()
                # Truncate for chat if too long
                preview = content[:500] + "..." if len(content) > 500 else content
                return f"Opening {found_file}:\n\n{preview}"
        except OSError as e:
            return f"I found the file but couldn't read it: {e}"
            
    return f"I couldn't find a document matching '{target}'."


def cmd_list_docs(args):
    """Usage: list documents / essays"""
    _ensure_docs_dir()
    
    try:
        files = os.listdir(DOCS_DIR)
    except OSError:
        return "I couldn't access the documents folder."
    
    if not files:
        return "The documents folder is empty."
    
    msg = "My Documents:\n"
    for f in files:
        msg += f" {f}\n"
    return msg


def cmd_create_ppt(args):
    """
    Creates a professional PPT. 
    Usage: automate PPT with slides: [ {title, content}, ... ]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"
    
    import json
    
    try:
        json_match = re.search(r'(\{.*\})', args, re.DOTALL)
        if not json_match:
            return "ERROR: Provide slide data in JSON format: { \"filename\": \"...\", \"slides\": [...] }"
            
        data = json.loads(json_match.group(1))
        filename = data.get("filename", "Presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        
        slides = data.get("slides", [])
        if not slides:
            return "ERROR: No slides provided in the JSON data."
        
        _ensure_docs_dir()
        filepath = os.path.join(DOCS_DIR, filename)
        
        prs = Presentation()
        
        for slide_data in slides:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_data.get("title", "Untitled Slide")
            body_shape.text = slide_data.get("content", "")
            
        prs.save(filepath)
        return f"SUCCESS: Professional PPT '{filename}' created with {len(slides)} slides and saved to documents folder. ✅"
        
    except json.JSONDecodeError as e:
        return f"ERROR [INVALID_JSON]: Could not parse slide data: {e}"
    except Exception as e:
        return f"ERROR [PPT_FAILURE]: {e}"


def cmd_create_pdf(args):
    """
    Creates a professional PDF report.
    Usage: automate PDF with title, author, content
    """
    try:
        from fpdf import FPDF
    except ImportError:
        return "ERROR: fpdf is not installed. Run: pip install fpdf"
    
    import json
    
    try:
        json_match = re.search(r'(\{.*\})', args, re.DOTALL)
        if not json_match:
            return "ERROR: Provide PDF data in JSON: { \"filename\": \"...\", \"title\": \"...\", \"content\": \"...\" }"
            
        data = json.loads(json_match.group(1))
        filename = data.get("filename", "Report.pdf")
        if not filename.endswith(".pdf"):
            filename += ".pdf"
        
        _ensure_docs_dir()
        filepath = os.path.join(DOCS_DIR, filename)
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(40, 10, data.get("title", "Academic Report"))
        pdf.ln(10)
        
        pdf.set_font("Arial", "", 12)
        pdf.multi_cell(0, 10, data.get("content", ""))
        
        pdf.output(filepath)
        return f"SUCCESS: Professional PDF '{filename}' created and saved to documents folder. ✅"
        
    except json.JSONDecodeError as e:
        return f"ERROR [INVALID_JSON]: Could not parse PDF data: {e}"
    except Exception as e:
        return f"ERROR [PDF_FAILURE]: {e}"


def register(dispatcher, nova_instance=None):
    global NOVA_CONTEXT
    if nova_instance:
        NOVA_CONTEXT = nova_instance
        
    dispatcher.register("write essay", cmd_write_essay)
    dispatcher.register("write an essay", cmd_write_essay)
    dispatcher.register("read essay", cmd_read_essay)
    dispatcher.register("access essay", cmd_read_essay)
    dispatcher.register("list essays", cmd_list_docs)
    dispatcher.register("list documents", cmd_list_docs)
    
    dispatcher.register("write email", cmd_write_email)
    dispatcher.register("write an email", cmd_write_email)
    dispatcher.register("generate email", cmd_write_email)
    dispatcher.register("write message", cmd_write_email)
    dispatcher.register("write a message", cmd_write_email)
    dispatcher.register("generate message", cmd_write_email)
    
    # Advanced Document Creators
    dispatcher.register("create ppt", cmd_create_ppt)
    dispatcher.register("create powerpoint", cmd_create_ppt)
    dispatcher.register("create pdf", cmd_create_pdf)
    dispatcher.register("generate report", cmd_create_pdf)
